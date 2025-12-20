"""
Industrial Procedural Knowledge Extraction - Document Processor
"""

import time
import hashlib
import asyncio
from typing import Dict, List, Any, Optional, Callable, TypeVar
from dataclasses import dataclass
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor
from threading import Lock

# Document processing libraries (optional)
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:  # pragma: no cover
    PANDAS_AVAILABLE = False
    pd = None

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_AVAILABLE = False
    Presentation = None  # type: ignore[assignment]

try:
    import fitz  # PyMuPDF
    FITZ_AVAILABLE = True
except ImportError:  # pragma: no cover
    FITZ_AVAILABLE = False
    fitz = None  # type: ignore[assignment]

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:  # pragma: no cover
    DOCX_AVAILABLE = False
    DocxDocument = None  # type: ignore[assignment]

# Audio processing
try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False
    whisper = None

# OCR
try:
    import easyocr
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    easyocr = None

from src.logging_config import get_logger
from src.core.unified_config import UnifiedConfig, get_config
from src.ai.types import ExtractionResult
from src.exceptions import ProcessingError

logger = get_logger(__name__)
T = TypeVar("T")


@dataclass
class ProcessingResult:
    """Simplified processing result"""
    document_id: str
    document_type: str
    content_extracted: str
    extraction_result: ExtractionResult
    processing_time: float
    file_size: int
    metadata: Dict[str, Any]


class StreamlinedDocumentProcessor:
    """
    Streamlined Document Processor
    
    Simplified, focused document processor that:
    1. Extracts content from various formats
    2. Uses the unified knowledge engine for extraction
    3. Provides consistent, fast processing
    """
    
    def __init__(self, config: Optional[UnifiedConfig] = None):
        self.config = config or get_config()
        self._engine: Optional["UnifiedKnowledgeEngine"] = None
        self.executor = ThreadPoolExecutor(max_workers=self.config.max_workers)
        self._ocr_reader = None
        self._ocr_lock = Lock()
        self._whisper_model = None
        self._whisper_lock = Lock()
        
        # Performance tracking
        self.stats = {
            'documents_processed': 0,
            'total_processing_time': 0.0,
            'format_counts': {},
            'average_processing_time': 0.0
        }
        
        # Supported formats
        self.text_formats = {'.pdf', '.doc', '.docx', '.txt', '.rtf'}
        self.image_formats = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'}
        self.spreadsheet_formats = {'.xls', '.xlsx', '.csv'}
        self.presentation_formats = {'.ppt', '.pptx'}
        self.audio_formats = {'.mp3', '.wav', '.flac', '.aac'}
        
        logger.info("Streamlined Document Processor initialized")

    @property
    def knowledge_engine(self) -> "UnifiedKnowledgeEngine":
        if self._engine is None:
            from src.ai.knowledge_engine import UnifiedKnowledgeEngine
            self._engine = UnifiedKnowledgeEngine(self.config)
        return self._engine
    
    async def _run_in_executor(self, func: Callable[[], T]) -> T:
        """Execute blocking operations in the shared thread pool."""
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(self.executor, func)
    
    async def process_document(
        self, 
        file_path: str, 
        document_id: Optional[str] = None
    ) -> ProcessingResult:
        """
        Process a document and extract knowledge
        
        Args:
            file_path: Path to the document file
            document_id: Optional document identifier
        """
        start_time = time.time()
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise ProcessingError(f"File not found: {file_path}")
        
        # Generate document ID if not provided
        if not document_id:
            document_id = hashlib.md5(str(file_path).encode()).hexdigest()[:16]
        
        # Detect file format and extract content
        file_format = file_path.suffix.lower()
        document_type = self._detect_document_type(file_format)
        
        try:
            content = await self._extract_content(file_path, file_format)
            
            if not content.strip():
                raise ProcessingError(f"No content extracted from {file_path}")
            
            # Extract knowledge using unified engine
            extraction_result = await self.knowledge_engine.extract_knowledge(
                content=content,
                document_type=document_type,
                quality_threshold=self.config.quality_threshold,
                document_id=document_id,
            )
            
            processing_time = time.time() - start_time
            file_size = file_path.stat().st_size
            
            # Update statistics
            self._update_stats(processing_time, file_format)
            
            result = ProcessingResult(
                document_id=document_id,
                document_type=document_type,
                content_extracted=content[:1000] + "..." if len(content) > 1000 else content,
                extraction_result=extraction_result,
                processing_time=processing_time,
                file_size=file_size,
                metadata={
                    'file_format': file_format,
                    'file_name': file_path.name,
                    'content_length': len(content),
                    'entities_extracted': len(extraction_result.entities),
                    'steps_extracted': len(extraction_result.steps),
                    'constraints_extracted': len(extraction_result.constraints),
                    'quality_metrics': extraction_result.quality_metrics
                }
            )
            
            logger.info(
                f"Processed {file_path.name} in {processing_time:.2f}s: "
                f"{len(extraction_result.steps)} steps, {len(extraction_result.entities)} entities extracted"
            )
            
            return result
            
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            raise ProcessingError(f"Processing failed: {e}") from e
    
    async def _extract_content(self, file_path: Path, file_format: str) -> str:
        """Extract text content based on file format"""
        
        if file_format in self.text_formats:
            return await self._extract_text_content(file_path, file_format)
        elif file_format in self.image_formats:
            return await self._extract_image_content(file_path)
        elif file_format in self.spreadsheet_formats:
            return await self._extract_spreadsheet_content(file_path, file_format)
        elif file_format in self.presentation_formats:
            return await self._extract_presentation_content(file_path)
        elif file_format in self.audio_formats:
            return await self._extract_audio_content(file_path)
        else:
            # Try to read as plain text
            return await self._extract_plain_text(file_path)
    
    async def _extract_text_content(self, file_path: Path, file_format: str) -> str:
        """Extract content from text documents"""
        
        if file_format == '.pdf':
            return await self._extract_pdf_content(file_path)
        elif file_format in {'.doc', '.docx'}:
            return await self._extract_docx_content(file_path)
        else:
            return await self._extract_plain_text(file_path)
    
    async def _extract_pdf_content(self, file_path: Path) -> str:
        """Extract text from PDF files"""
        if not FITZ_AVAILABLE:
            logger.warning("PyMuPDF not available. Install PyMuPDF to process PDFs.")
            return ""

        def extract_pdf():
            try:
                doc = fitz.open(str(file_path))
                try:
                    pages = [page.get_text() for page in doc]
                    return "".join(pages)
                finally:
                    doc.close()
            except Exception as e:
                logger.warning(f"PDF extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_pdf)
    
    async def _extract_docx_content(self, file_path: Path) -> str:
        """Extract text from Word documents"""
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not available. Install python-docx to process Word documents.")
            return ""

        def extract_docx():
            try:
                doc = DocxDocument(str(file_path))
                paragraphs = [paragraph.text for paragraph in doc.paragraphs]
                return "\n".join(paragraphs)
            except Exception as e:
                logger.warning(f"DOCX extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_docx)
    
    def _get_ocr_reader(self):
        if self._ocr_reader is not None:
            return self._ocr_reader
        with self._ocr_lock:
            if self._ocr_reader is None and OCR_AVAILABLE:
                self._ocr_reader = easyocr.Reader(['en'])
        return self._ocr_reader

    async def _extract_image_content(self, file_path: Path) -> str:
        """Extract text from images using OCR"""
        if not OCR_AVAILABLE:
            logger.warning("OCR not available. Install easyocr for image processing.")
            return ""
        
        def extract_ocr():
            try:
                reader = self._get_ocr_reader()
                if reader is None:
                    return ""
                with self._ocr_lock:
                    results = reader.readtext(str(file_path))
                return " ".join(result[1] for result in results)
            except Exception as e:
                logger.warning(f"OCR extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_ocr)
    
    async def _extract_spreadsheet_content(self, file_path: Path, file_format: str) -> str:
        """Extract content from spreadsheets"""
        if not PANDAS_AVAILABLE:
            logger.warning("pandas not available. Install pandas to process spreadsheet documents.")
            return ""

        def extract_spreadsheet():
            try:
                if file_format == '.csv':
                    df = pd.read_csv(file_path)
                else:
                    df = pd.read_excel(file_path)
                
                # Convert to text representation
                text = f"Spreadsheet data ({df.shape[0]} rows, {df.shape[1]} columns):\n"
                text += df.to_string(max_rows=50, max_cols=10)
                return text
            except Exception as e:
                logger.warning(f"Spreadsheet extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_spreadsheet)
    
    async def _extract_presentation_content(self, file_path: Path) -> str:
        """Extract text from presentations"""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available. Install python-pptx to process PowerPoint files.")
            return ""

        def extract_presentation():
            try:
                prs = Presentation(str(file_path))
                fragments = []

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            fragments.append(shape.text)

                return "\n".join(fragments)
            except Exception as e:
                logger.warning(f"Presentation extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_presentation)
    
    def _get_whisper_model(self):
        if self._whisper_model is not None:
            return self._whisper_model
        with self._whisper_lock:
            if self._whisper_model is None and WHISPER_AVAILABLE:
                model_id = getattr(self.config, "whisper_model", "base")
                self._whisper_model = whisper.load_model(model_id)
        return self._whisper_model

    async def _extract_audio_content(self, file_path: Path) -> str:
        """Extract text from audio files using Whisper"""
        if not WHISPER_AVAILABLE:
            logger.warning("Whisper not available. Install openai-whisper for audio processing.")
            return ""
        
        def extract_audio():
            try:
                model = self._get_whisper_model()
                if model is None:
                    return ""
                with self._whisper_lock:
                    result = model.transcribe(str(file_path))
                return result.get("text", "")
            except Exception as e:
                logger.warning(f"Audio extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_audio)
    
    async def _extract_plain_text(self, file_path: Path) -> str:
        """Extract plain text content"""
        def extract_text():
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    return file.read()
            except Exception as e:
                logger.warning(f"Plain text extraction failed: {e}")
                return ""
        
        return await self._run_in_executor(extract_text)
    
    def _detect_document_type(self, file_format: str) -> str:
        """Detect document type based on file format"""
        type_mapping = {
            '.pdf': 'manual',
            '.doc': 'document',
            '.docx': 'document',
            '.txt': 'text',
            '.ppt': 'presentation',
            '.pptx': 'presentation',
            '.xls': 'spreadsheet',
            '.xlsx': 'spreadsheet',
            '.csv': 'data',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.png': 'image',
            '.mp3': 'audio',
            '.wav': 'audio',
        }
        return type_mapping.get(file_format, 'unknown')
    
    def _update_stats(self, processing_time: float, file_format: str):
        """Update processing statistics"""
        self.stats['documents_processed'] += 1
        self.stats['total_processing_time'] += processing_time
        self.stats['format_counts'][file_format] = self.stats['format_counts'].get(file_format, 0) + 1
        self.stats['average_processing_time'] = (
            self.stats['total_processing_time'] / self.stats['documents_processed']
        )
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        stats = self.stats.copy()
        stats['knowledge_engine_stats'] = self.knowledge_engine.get_performance_stats()
        return stats
    
    def clear_cache(self):
        """Clear all caches"""
        self.knowledge_engine.clear_cache()

    def shutdown(self) -> None:
        """Release processor resources."""
        self.executor.shutdown(wait=False)
        if self._engine is not None:
            self._engine.clear_cache()


# Backward compatibility
OptimizedDocumentProcessor = StreamlinedDocumentProcessor


# Factory function
def create_document_processor(config: Optional[UnifiedConfig] = None) -> StreamlinedDocumentProcessor:
    """Create a new streamlined document processor"""
    return StreamlinedDocumentProcessor(config)
