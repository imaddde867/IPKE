#!/usr/bin/env python3
"""
OPTIMIZED Document Processor - Performance-First Implementation
Target: 2 minutes max per document processing time
Current: 10+ minutes per document. Improvement needed: 5x speed increase minimum.

This processor is designed for SPEED FIRST, then quality, optimized for M4 MacBook Pro.
"""

import asyncio
import logging
import time
import json
import hashlib
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Callable
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor
import threading
import os
import io
import tempfile
from datetime import datetime

# Core processing libraries (lazy loading)
import requests
import cv2
import numpy as np
from PIL import Image, ImageEnhance
# import pytesseract  # Replaced with EasyOCR
import pandas as pd
from pptx import Presentation
import PyPDF2
import fitz  # PyMuPDF
from docx import Document as DocxDocument

# AI and NLP libraries (lazy loading)
import spacy
try:
    import torch
    TORCH_AVAILABLE = True
except ImportError:
    TORCH_AVAILABLE = False
    torch = None

# Internal imports
from src.logging_config import get_logger, log_processing_step
from src.core.config import config as config_manager
from src.exceptions import ProcessingError, AIError
from src.ai.advanced_knowledge_engine import AdvancedKnowledgeEngine
from src.ai.llm_processing_engine import LLMProcessingEngine
from src.ai.enhanced_extraction_engine import EnhancedExtractionEngine

logger = get_logger(__name__)

@dataclass
class ProcessingResult:
    """Optimized processing result"""
    document_id: str
    document_type: str
    processing_time: float
    entities_extracted: int
    confidence_score: float
    performance_metrics: Dict[str, Any]
    entities: List[Dict[str, Any]]
    processing_method: str
    cache_hit: bool
    optimization_level: str
    content_summary: str
    file_size: int
    format_detected: str

class OptimizedDocumentProcessor:
    """OPTIMIZED document processor with parallel processing and caching"""
    
    def __init__(self, config: Dict[str, Any] = None, db_session=None):
        self.config = config or {}
        self.db_session = db_session
        
        # Performance optimizations
        self.max_workers = 4  # Optimize for M4 chip
        self.executor = ThreadPoolExecutor(max_workers=self.max_workers)
        self.processing_cache = {}
        self.cache_lock = threading.Lock()
        
        # Initialize REAL AI engines (lazy loading)
        self.advanced_engine = None
        self.llm_engine = None
        self.extraction_engine = None
        self.engines_initialized = False
        
        # Performance monitoring
        self.processing_stats = {
            "total_documents": 0,
            "total_processing_time": 0.0,
            "cache_hits": 0,
            "cache_misses": 0,
            "average_processing_time": 0.0,
            "performance_target_met": 0,
            "performance_target_missed": 0,
            "format_processing_times": {},
            "error_count": 0
        }
        
        # Performance thresholds
        self.target_processing_time = 120.0  # 2 minutes target
        self.performance_warning_threshold = 90.0  # 1.5 minutes warning
        
        # Supported formats (optimized for speed)
        self.supported_formats = {
            'text': ['.pdf', '.doc', '.docx', '.txt', '.rtf'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
            'spreadsheet': ['.xls', '.xlsx', '.csv'],
            'presentation': ['.ppt', '.pptx'],
            'audio': ['.mp3', '.wav', '.flac', '.aac'],
            'video': ['.mp4', '.avi', '.mov', '.mkv']
        }
        
        # Initialize basic capabilities
        self._init_basic_capabilities()
        
        print("🚀 OPTIMIZED Document Processor initialized")
        print(f"🎯 Performance Target: {self.target_processing_time} seconds per document")
        print(f"⚡ Max Workers: {self.max_workers} (optimized for M4)")
        print(f"📁 Supported Formats: {len(self.supported_formats)} categories")

    # Backward-compatibility: legacy callers expect process_document()
    def process_document(self, file_path: str, document_id: str = None) -> Dict[str, Any]:
        """Legacy sync entrypoint that returns a dict shaped for the frontend converter.

        This wraps process_document_sync and adapts the result to the
        structure expected by convert_document_processor_results_to_display().
        """
        result = self.process_document_sync(file_path, document_id)

        # Determine method labels compatible with the frontend expectations
        primary_method = result.processing_method
        is_llm_primary = primary_method in {"ai_enhanced_processing", "ai_processing"}
        processing_meta = {
            "method": "llm_first_processing" if is_llm_primary else "advanced_engine",
            "processing_method": "llm_primary" if is_llm_primary else primary_method,
            "confidence_score": result.confidence_score,
            "processing_time": result.processing_time,
            "validation_passed": True
        }

        # Normalize entities to the expected schema
        extracted_entities: List[Dict[str, Any]] = []
        for ent in result.entities:
            if not isinstance(ent, dict):
                continue
            core_content = ent.get("core_content") or ent.get("content") or ent.get("key_identifier") or ""
            extracted_entities.append({
                "core_content": core_content,
                "key_identifier": ent.get("key_identifier", (core_content or "")[:50]),
                "category": ent.get("category", ent.get("entity_type", "unknown")),
                "entity_type": ent.get("entity_type", ent.get("category", "unknown")),
                "confidence_score": ent.get("confidence", result.confidence_score),
                "source_section": ent.get("source_location", ""),
                "priority_level": "high" if ent.get("confidence", 0.0) >= 0.85 else ("medium" if ent.get("confidence", 0.0) >= 0.6 else "low"),
                "context_tags": ent.get("relationships", []),
                "completeness_score": ent.get("completeness", 0.8),
                "clarity_score": ent.get("clarity", 0.8),
                "actionability_score": ent.get("actionability", 0.7)
            })

        return {
            "knowledge": {
                "extracted_entities": extracted_entities,
                "processing_metadata": processing_meta,
                "extraction_methods": [result.processing_method],
            },
            "summary": result.content_summary,
            "file": Path(file_path).name
        }
    
    def _init_basic_capabilities(self):
        """Initialize basic processing capabilities"""
        # Test EasyOCR availability (primary OCR method) without signal-based timeout
        try:
            import easyocr
            import threading
            import time
            
            # Use threading-based timeout instead of signal (works in Streamlit)
            def test_easyocr():
                try:
                    test_reader = easyocr.Reader(['en'], gpu=False)
                    return True, None
                except Exception as e:
                    return False, e
            
            # Start EasyOCR test in a separate thread with timeout
            result_container = [None]
            error_container = [None]
            
            def worker():
                try:
                    success, error = test_easyocr()
                    result_container[0] = success
                    error_container[0] = error
                except Exception as e:
                    result_container[0] = False
                    error_container[0] = e
            
            thread = threading.Thread(target=worker)
            thread.daemon = True
            thread.start()
            thread.join(timeout=10)  # 10 second timeout
            
            if thread.is_alive():
                # Thread is still running, timeout occurred
                logger.warning("EasyOCR initialization timed out")
                self.ocr_available = False
            else:
                # Thread completed
                if result_container[0]:
                    self.ocr_available = True
                    logger.info("EasyOCR initialized successfully")
                else:
                    self.ocr_available = False
                    logger.warning(f"EasyOCR initialization failed: {error_container[0]}")
                    
        except Exception as e:
            logger.warning(f"EasyOCR initialization failed: {e}")
            self.ocr_available = False
        
        # Initialize other capabilities as needed
        self.audio_processing_available = False
        self.video_processing_available = False
        
        # Test basic libraries
        try:
            import whisper
            self.audio_processing_available = True
        except ImportError:
            pass
        
        try:
            import cv2
            self.video_processing_available = True
        except ImportError:
            pass
    
    def _initialize_engines(self):
        """Lazy initialize processing engines - skip for HPC compatibility"""
        if self.engines_initialized:
            return
        
        # Skip AI engine initialization on HPC systems to prevent hanging
        print("⚠️ Skipping AI engine initialization for HPC compatibility")
        print("ℹ️ AI engines will be initialized when actually needed")
        
        self.advanced_engine = None
        self.llm_engine = None
        self.extraction_engine = None
        self.engines_initialized = True
        
        # Note: AI engines will be initialized lazily when process_document is called
    
    def _initialize_engines_lazy(self):
        """Initialize AI engines only when actually needed for processing"""
        if self.engines_initialized and (self.advanced_engine or self.llm_engine or self.extraction_engine):
            return  # Already initialized
        
        print("🔄 Initializing AI engines for document processing...")
        
        try:
            import os
            
            # Set offline mode to prevent hanging
            os.environ['TRANSFORMERS_OFFLINE'] = '1'
            os.environ['HF_HUB_OFFLINE'] = '1'
            
            # Use threading-based timeout for AI engine initialization (works in Streamlit)
            import threading
            import time
            
            result_container = [None]
            error_container = [None]
            
            def engine_init_worker():
                try:
                    # Initialize advanced knowledge engine
                    if not self.advanced_engine:
                        self.advanced_engine = AdvancedKnowledgeEngine(config_manager.ai, self.db_session)
                        print("✅ Advanced Knowledge Engine initialized")
                    
                    # Initialize LLM processing engine
                    if not self.llm_engine:
                        self.llm_engine = LLMProcessingEngine()
                        # Initialize the LLM engine asynchronously
                        import asyncio
                        try:
                            # Create a new event loop for this thread
                            loop = asyncio.new_event_loop()
                            asyncio.set_event_loop(loop)
                            loop.run_until_complete(self.llm_engine.initialize())
                            loop.close()
                            print("✅ LLM Processing Engine initialized")
                        except Exception as e:
                            print(f"⚠️ LLM engine initialization failed: {e}")
                            self.llm_engine = None
                    
                    # Initialize enhanced extraction engine
                    if not self.extraction_engine:
                        self.extraction_engine = EnhancedExtractionEngine()
                        print("✅ Enhanced Extraction Engine initialized")
                    
                    result_container[0] = True
                    
                except Exception as e:
                    error_container[0] = e
            
            # Start AI engine initialization in a separate thread with timeout
            thread = threading.Thread(target=engine_init_worker)
            thread.daemon = True
            thread.start()
            thread.join(timeout=10)  # 10 second timeout for AI engine initialization
            
            if thread.is_alive():
                # Thread is still running, timeout occurred
                print("⚠️ AI engine initialization timed out - continuing without AI engines")
                self.advanced_engine = None
                self.llm_engine = None
                self.extraction_engine = None
            else:
                # Thread completed
                if error_container[0]:
                    print(f"⚠️ Engine initialization failed: {error_container[0]}")
                    self.advanced_engine = None
                    self.llm_engine = None
                    self.extraction_engine = None
                else:
                    print("✅ All AI engines initialized successfully")
            
            self.engines_initialized = True
                
        except Exception as e:
            print(f"⚠️ Engine initialization failed: {e}")
            self.advanced_engine = None
            self.llm_engine = None
            self.extraction_engine = None
            self.engines_initialized = True
    
    async def _async_initialize_engines(self):
        """Async initialization of engines"""
        try:
            if self.advanced_engine:
                await self.advanced_engine.initialize()
                print("✅ Advanced Knowledge Engine async initialization completed")
            
            if self.llm_engine:
                await self.llm_engine.initialize()
                print("✅ LLM Processing Engine async initialization completed")
                
        except Exception as e:
            print(f"⚠️ Async engine initialization failed: {e}")
            logger.error(f"Async engine initialization failed: {e}")
    
    async def process_document_async(self, file_path: str, document_id: str = None) -> ProcessingResult:
        """OPTIMIZED async document processing"""
        start_time = time.time()
        
        # Generate document ID if not provided
        if not document_id:
            file_path_obj = Path(file_path)
            content_hash = hashlib.md5(f"{file_path}_{start_time}".encode()).hexdigest()
            document_id = f"doc_{file_path_obj.stem}_{content_hash[:8]}"
        
        # Check cache first
        cache_key = self._generate_cache_key(file_path)
        cached_result = self._get_cached_result(cache_key)
        
        if cached_result:
            self.processing_stats["cache_hits"] += 1
            return cached_result
        
        self.processing_stats["cache_misses"] += 1
        
        try:
            # Get file info
            file_path_obj = Path(file_path)
            file_size = file_path_obj.stat().st_size
            file_extension = file_path_obj.suffix.lower()
            format_detected = self._get_file_type(file_extension)
            
            # Extract content based on format
            content = await self._extract_content_async(file_path, format_detected)
            
            # Initialize engines if needed
            self._initialize_engines_lazy()
            
            # Parallel processing pipeline
            processing_tasks = [
                self._process_with_llm_async(content, format_detected),
                self._extract_entities_async(content, format_detected),
                self._validate_and_enhance_async(content, format_detected)
            ]
            
            # Execute all tasks in parallel
            try:
                results = await asyncio.gather(*processing_tasks, return_exceptions=True)
                
                # Process results
                llm_result = results[0] if not isinstance(results[0], Exception) else None
                extraction_result = results[1] if not isinstance(results[1], Exception) else None
                validation_result = results[2] if not isinstance(results[2], Exception) else None
                
            except Exception as e:
                logger.error(f"Parallel processing failed: {e}")
                # Fallback to sequential processing
                llm_result = await self._process_with_llm_async(content, format_detected)
                extraction_result = await self._extract_entities_async(content, format_detected)
                validation_result = await self._validate_and_enhance_async(content, format_detected)
            
            # Merge and finalize results
            final_result = self._merge_processing_results(
                llm_result, extraction_result, validation_result, 
                content, format_detected, document_id, file_size
            )
            
            # Calculate processing time
            processing_time = time.time() - start_time
            final_result.processing_time = processing_time
            
            # Update performance stats
            self._update_performance_stats(processing_time, format_detected)
            
            # Cache the result
            self._cache_result(cache_key, final_result)
            
            # Performance monitoring
            self._monitor_performance(processing_time, document_id)
            
            return final_result
            
        except Exception as e:
            logger.error(f"Document processing failed: {e}")
            self.processing_stats["error_count"] += 1
            
            # Return error result
            return ProcessingResult(
                document_id=document_id,
                document_type="error",
                processing_time=time.time() - start_time,
                entities_extracted=0,
                confidence_score=0.0,
                performance_metrics={"error": str(e)},
                entities=[],
                processing_method="error",
                cache_hit=False,
                optimization_level="error",
                content_summary=f"Processing failed: {str(e)}",
                file_size=0,
                format_detected="unknown"
            )
    
    def process_document_sync(self, file_path: str, document_id: str = None) -> ProcessingResult:
        """Synchronous wrapper for async processing"""
        try:
            return asyncio.run(self.process_document_async(file_path, document_id))
        except RuntimeError:
            # If already in event loop, create new one
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                return loop.run_until_complete(
                    self.process_document_async(file_path, document_id)
                )
            finally:
                loop.close()
    
    async def _extract_content_async(self, file_path: str, format_type: str) -> str:
        """Async content extraction based on format"""
        try:
            # Use ThreadPoolExecutor for I/O-bound operations
            loop = asyncio.get_event_loop()
            
            if format_type == "text":
                return await loop.run_in_executor(
                    self.executor, self._extract_text_content, file_path
                )
            elif format_type == "image":
                return await loop.run_in_executor(
                    self.executor, self._extract_image_content, file_path
                )
            elif format_type == "spreadsheet":
                return await loop.run_in_executor(
                    self.executor, self._extract_spreadsheet_content, file_path
                )
            elif format_type == "presentation":
                return await loop.run_in_executor(
                    self.executor, self._extract_presentation_content, file_path
                )
            elif format_type == "audio":
                return await loop.run_in_executor(
                    self.executor, self._extract_audio_content, file_path
                )
            elif format_type == "video":
                return await loop.run_in_executor(
                    self.executor, self._extract_video_content, file_path
                )
            else:
                return await loop.run_in_executor(
                    self.executor, self._extract_generic_content, file_path
                )
                
        except Exception as e:
            logger.error(f"Content extraction failed: {e}")
            return f"Content extraction failed: {str(e)}"
    
    def _extract_text_content(self, file_path: str) -> str:
        """Extract text content from text-based documents"""
        try:
            file_path_obj = Path(file_path)
            extension = file_path_obj.suffix.lower()
            
            if extension == '.pdf':
                return self._extract_pdf_content(file_path_obj)
            elif extension in ['.doc', '.docx']:
                return self._extract_word_content(file_path_obj)
            elif extension in ['.txt', '.rtf']:
                return self._extract_text_content_simple(file_path_obj)
            else:
                return self._extract_generic_text(file_path_obj)
                
        except Exception as e:
            logger.error(f"Text extraction failed: {e}")
            return f"Text extraction failed: {str(e)}"
    
    def _extract_pdf_content(self, file_path: Path) -> str:
        """Fast PDF content extraction with OCR fallback for image-based PDFs"""
        try:
            # Try PyMuPDF first (faster)
            doc = fitz.open(file_path)
            text = ""
            page_count = 0
            for page in doc:
                page_text = page.get_text()
                text += page_text
                page_count += 1
                print(f"📄 PDF Page {page_count}: {len(page_text)} characters")
            doc.close()
            
            print(f"📊 PDF Total: {len(text)} characters from {page_count} pages")
            
            # If text extraction failed (image-based PDF), try OCR
            if len(text.strip()) < 50:
                print(f"⚠️ PDF appears to be image-based (only {len(text)} chars)")
                print(f"🔄 Quick OCR check on first 3 pages...")
                
                # Try OCR processing for image-based PDF - only check first 3 pages for speed
                ocr_text = self._extract_pdf_with_ocr(file_path, 3)  # Only check first 3 pages
                if ocr_text and len(ocr_text.strip()) > 50:
                    print(f"✅ OCR successful: {len(ocr_text)} characters extracted")
                    return ocr_text
                else:
                    print(f"❌ OCR also failed or returned minimal content")
                    return f"PDF appears to be image-based but OCR failed. Text chars: {len(text)}, OCR chars: {len(ocr_text) if ocr_text else 0}"
            
            return text
        except Exception as e:
            print(f"❌ PyMuPDF failed: {e}")
            try:
                # Fallback to PyPDF2
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    page_count = len(reader.pages)
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text()
                        text += page_text
                        print(f"📄 PDF Page {i+1}: {len(page_text)} characters")
                    
                    print(f"📊 PDF Total (PyPDF2): {len(text)} characters from {page_count} pages")
                    
                    # If PyPDF2 also fails, try OCR
                    if len(text.strip()) < 50:
                        print(f"⚠️ PyPDF2 also found minimal text, trying OCR...")
                        ocr_text = self._extract_pdf_with_ocr(file_path, page_count)
                        if ocr_text and len(ocr_text.strip()) > 50:
                            print(f"✅ OCR successful: {len(ocr_text)} characters extracted")
                            return ocr_text
                        else:
                            return f"PDF appears to be image-based but OCR failed. Text chars: {len(text)}, OCR chars: {len(ocr_text) if ocr_text else 0}"
                    
                    return text
            except Exception as e2:
                print(f"❌ PyPDF2 also failed: {e2}")
                return f"PDF extraction failed: PyMuPDF: {e}, PyPDF2: {e2}"
    
    def _extract_pdf_with_ocr(self, file_path: Path, max_pages: int = 10) -> str:
        """Extract text from image-based PDF using OCR"""
        if not self.ocr_available:
            print(f"❌ OCR not available for image-based PDF processing")
            return ""
        
        try:
            import fitz  # PyMuPDF for PDF to image conversion
            import easyocr
            import cv2
            import numpy as np
            
            print(f"🔍 Starting OCR processing for image-based PDF...")
            
            # Create EasyOCR reader once and reuse for all pages
            reader = easyocr.Reader(['en'], gpu=False)
            
            # Convert PDF pages to images and process with OCR
            doc = fitz.open(file_path)
            ocr_text = ""
            pages_processed = 0
            
            # Limit to first 10 pages for performance (can be adjusted)
            max_pages_to_process = min(max_pages, 10)
            
            for page_num in range(max_pages_to_process):
                try:
                    print(f"🔍 Processing page {page_num + 1}...")
                    page = doc[page_num]
                    # Convert page to image
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")
                    
                    # Convert to OpenCV format
                    nparr = np.frombuffer(img_data, np.uint8)
                    img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                    
                    print(f"📷 Image size: {img.shape if img is not None else 'None'}")
                    
                    if img is not None:
                        # Preprocess image for better OCR
                        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                        enhanced = cv2.equalizeHist(gray)
                        
                        # Use EasyOCR (reader already created above)
                        result = reader.readtext(enhanced)
                        
                        # Extract text from OCR results
                        page_text = ""
                        for (bbox, detected_text, confidence) in result:
                            if confidence > 0.3:  # Lower threshold for better coverage
                                page_text += detected_text + " "
                        
                        if page_text.strip():
                            ocr_text += f"\n--- Page {page_num + 1} ---\n{page_text.strip()}\n"
                            pages_processed += 1
                            print(f"📄 OCR Page {page_num + 1}: {len(page_text)} characters")
                        
                except Exception as page_error:
                    print(f"⚠️ OCR failed for page {page_num + 1}: {page_error}")
                    continue
            
            doc.close()
            
            print(f"📊 OCR Total: {len(ocr_text)} characters from {pages_processed} pages")
            return ocr_text
            
        except Exception as e:
            print(f"❌ OCR processing failed: {e}")
            return ""
    
    def _extract_word_content(self, file_path: Path) -> str:
        """Fast Word document content extraction"""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            return f"Word extraction failed: {str(e)}"
    
    def _extract_text_content_simple(self, file_path: Path) -> str:
        """Simple text file extraction"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return f"Text extraction failed: {str(e)}"
    
    def _extract_generic_text(self, file_path: Path) -> str:
        """Generic text extraction fallback"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            return f"Generic extraction failed: {str(e)}"
    
    def _extract_image_content(self, file_path: str) -> str:
        """Enhanced image content extraction with EasyOCR, preprocessing, and debugging"""
        try:
            if not self.ocr_available:
                return "OCR not available for image processing"
            
            import signal
            import easyocr
            import cv2
            import numpy as np
            from PIL import Image, ImageEnhance
            
            # Use threading-based timeout instead of signal (works in Streamlit)
            import threading
            import time
            
            # Preprocess image for better OCR
            print(f"🔍 Processing image: {Path(file_path).name}")
            
            # Load and preprocess image
            img = cv2.imread(file_path)
            if img is None:
                return "Could not load image file"
            
            # Convert to grayscale
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            
            # Enhance contrast
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
            enhanced = clahe.apply(gray)
            
            # Denoise
            denoised = cv2.fastNlMeansDenoising(enhanced)
            
            # Save preprocessed image temporarily
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                cv2.imwrite(temp_file.name, denoised)
                temp_path = temp_file.name
            
            try:
                # Use threading-based timeout for OCR processing
                result_container = [None]
                error_container = [None]
                
                def ocr_worker():
                    try:
                        # Use EasyOCR with multiple languages and lower confidence threshold
                        reader = easyocr.Reader(['en'], gpu=False)
                        
                        # Extract text from preprocessed image
                        result = reader.readtext(temp_path)
                        
                        print(f"📊 EasyOCR detected {len(result)} text regions:")
                        
                        # Process results with detailed debugging
                        text_parts = []
                        all_text_parts = []  # For debugging
                        
                        for i, (bbox, detected_text, confidence) in enumerate(result):
                            all_text_parts.append(f"Region {i+1}: '{detected_text}' (conf: {confidence:.2f})")
                            
                            # Lower confidence threshold for better text capture
                            if confidence > 0.3:  # Reduced from 0.5 to 0.3
                                text_parts.append(detected_text)
                                print(f"  ✅ Region {i+1}: '{detected_text}' (conf: {confidence:.2f})")
                            else:
                                print(f"  ❌ Region {i+1}: '{detected_text}' (conf: {confidence:.2f}) - filtered out")
                        
                        # Print all detected regions for debugging
                        if all_text_parts:
                            print("🔍 All detected text regions:")
                            for part in all_text_parts:
                                print(f"  {part}")
                        
                        combined_text = ' '.join(text_parts)
                        
                        if combined_text.strip():
                            print(f"✅ Final extracted text ({len(combined_text)} chars): '{combined_text[:100]}{'...' if len(combined_text) > 100 else ''}'")
                            logger.info(f"EasyOCR extracted {len(combined_text)} characters from image")
                            result_container[0] = combined_text
                        else:
                            print("❌ No text passed confidence threshold")
                            result_container[0] = "No text detected in image (all regions below confidence threshold)"
                            
                    except Exception as e:
                        error_container[0] = e
                
                # Start OCR in a separate thread with timeout
                thread = threading.Thread(target=ocr_worker)
                thread.daemon = True
                thread.start()
                thread.join(timeout=60)  # 60 second timeout for OCR
                
                if thread.is_alive():
                    # Thread is still running, timeout occurred
                    logger.warning("OCR processing timed out")
                    return "OCR processing timed out - image too complex or system overloaded"
                else:
                    # Thread completed
                    if error_container[0]:
                        return f"Image extraction failed: {str(error_container[0])}"
                    elif result_container[0]:
                        return result_container[0]
                    else:
                        return "No text detected in image"
                        
            finally:
                # Clean up temp file
                import os
                try:
                    os.unlink(temp_path)
                except Exception:
                    pass
            
        except Exception as e:
            print(f"❌ OCR processing failed: {e}")
            return f"Image extraction failed: {str(e)}"
    
    def _extract_spreadsheet_content(self, file_path: str) -> str:
        """Fast spreadsheet content extraction"""
        try:
            file_path_obj = Path(file_path)
            extension = file_path_obj.suffix.lower()
            
            if extension == '.csv':
                df = pd.read_csv(file_path)
            elif extension in ['.xls', '.xlsx']:
                df = pd.read_excel(file_path)
            else:
                return "Unsupported spreadsheet format"
            
            # Convert to text representation
            return df.to_string()
            
        except Exception as e:
            return f"Spreadsheet extraction failed: {str(e)}"
    
    def _extract_presentation_content(self, file_path: str) -> str:
        """Fast presentation content extraction"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text
            
        except Exception as e:
            return f"Presentation extraction failed: {str(e)}"
    
    def _extract_audio_content(self, file_path: str) -> str:
        """Fast audio content extraction"""
        if not self.audio_processing_available:
            return "Audio processing not available"
        
        try:
            # Basic audio metadata extraction
            return f"Audio file: {Path(file_path).name} (transcription not available)"
        except Exception as e:
            return f"Audio extraction failed: {str(e)}"
    
    def _extract_video_content(self, file_path: str) -> str:
        """Fast video content extraction via sparse frame OCR with EasyOCR"""
        if not self.video_processing_available:
            return "Video processing not available"
        if not self.ocr_available:
            return f"Video file: {Path(file_path).name} (OCR not available)"
        
        try:
            result = self._process_video_document(Path(file_path))
            text = (result.get('text') or '').strip()
            # If no text from frame OCR, try audio transcription using Whisper if available
            if not text and self.audio_processing_available:
                try:
                    import whisper  # type: ignore
                    model = whisper.load_model("base")
                    transcribe = model.transcribe(str(file_path), fp16=False)
                    audio_text = (transcribe.get("text") or "").strip()
                    if audio_text:
                        text = audio_text
                except Exception as _audio_err:
                    # Keep silent fallback; return whatever we have
                    pass
            if text:
                return text
            return f"Video file: {Path(file_path).name} (no readable text detected)"
        except Exception as e:
            return f"Video extraction failed: {str(e)}"
    
    def _extract_generic_content(self, file_path: str) -> str:
        """Generic content extraction fallback"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception:
            return f"Generic content extraction failed for {Path(file_path).name}"
    
    async def _process_with_llm_async(self, content: str, document_type: str) -> Dict[str, Any]:
        """Async LLM processing using real AI engine"""
        if not self.llm_engine:
            return {"entities": [], "confidence": 0.0, "method": "llm_unavailable"}
        
        try:
            # Debug logging for content received by LLM
            print(f"🧠 LLM Processing: {document_type} document")
            print(f"📝 Content length: {len(content) if content else 0} characters")
            if content and len(content) < 200:
                print(f"📄 Content preview: '{content[:200]}...'")
            
            # Use the real LLM engine's process_document method
            # Pass metadata to encourage LLM processing for video/image
            metadata = {
                'file_type': 'video' if document_type == 'video' else ('image' if document_type == 'image' else document_type),
                'content_type': 'video_content' if document_type == 'video' else 'document'
            }
            # Ensure non-empty content to avoid degenerate LLM calls
            safe_content = content if isinstance(content, str) and content.strip() else ""
            
            if not safe_content or len(safe_content.strip()) < 10:
                print(f"⚠️ LLM skipping: content too short ({len(safe_content)} chars)")
                return {"entities": [], "confidence": 0.0, "method": "content_too_short"}
            
            result = await self.llm_engine.process_document(safe_content, document_type, metadata)
            
            # Convert ProcessingResult to dict format
            if hasattr(result, 'entities'):
                entities_dict = []
                for entity in result.entities:
                    if hasattr(entity, '__dict__'):
                        entity_dict = entity.__dict__.copy()
                    else:
                        entity_dict = asdict(entity) if hasattr(entity, '__dict__') else str(entity)
                    entities_dict.append(entity_dict)
                
                return {
                    "entities": entities_dict,
                    "confidence": result.confidence_score if hasattr(result, 'confidence_score') else 0.0,
                    "method": result.processing_method if hasattr(result, 'processing_method') else "llm_processing",
                    "quality_metrics": result.quality_metrics if hasattr(result, 'quality_metrics') else {},
                    "llm_enhanced": result.llm_enhanced if hasattr(result, 'llm_enhanced') else True,
                    "validation_passed": result.validation_passed if hasattr(result, 'validation_passed') else True
                }
            else:
                # Fallback if result doesn't have expected structure
                return {
                    "entities": [],
                    "confidence": 0.8,
                    "method": "llm_processing_fallback",
                    "quality_metrics": {},
                    "llm_enhanced": True,
                    "validation_passed": True
                }
                
        except Exception as e:
            logger.error(f"LLM processing failed: {e}")
            return {"entities": [], "confidence": 0.0, "method": "llm_failed"}
    
    async def _extract_entities_async(self, content: str, document_type: str) -> Dict[str, Any]:
        """Async entity extraction using real AI engine"""
        if not self.extraction_engine:
            return {"entities": [], "confidence": 0.0, "method": "extraction_unavailable"}
        
        try:
            # Use the real extraction engine's extract_comprehensive_knowledge method
            result = self.extraction_engine.extract_comprehensive_knowledge(content, document_type)
            
            # Convert ExtractedEntity objects to dict format
            entities_dict = []
            total_confidence = 0.0
            entity_count = 0
            
            for entity in result:
                if hasattr(entity, 'content') and hasattr(entity, 'confidence'):
                    entity_dict = {
                        'content': entity.content,
                        'entity_type': getattr(entity, 'entity_type', 'unknown'),
                        'category': getattr(entity, 'category', 'general'),
                        'confidence': entity.confidence,
                        'context': getattr(entity, 'context', ''),
                        'metadata': getattr(entity, 'metadata', {}),
                        'relationships': getattr(entity, 'relationships', []),
                        'source_location': getattr(entity, 'source_location', '')
                    }
                    entities_dict.append(entity_dict)
                    total_confidence += entity.confidence
                    entity_count += 1
            
            avg_confidence = total_confidence / entity_count if entity_count > 0 else 0.0
            
            return {
                "entities": entities_dict,
                "confidence": avg_confidence,
                "method": "enhanced_extraction",
                "entity_count": entity_count,
                "extraction_quality": "high" if avg_confidence > 0.7 else "medium"
            }
            
        except Exception as e:
            logger.error(f"Entity extraction failed: {e}")
            return {"entities": [], "confidence": 0.0, "method": "extraction_failed"}
    
    async def _validate_and_enhance_async(self, content: str, document_type: str) -> Dict[str, Any]:
        """Async validation and enhancement using real AI engine"""
        try:
            if not self.advanced_engine:
                # Fallback to basic validation
                return self._basic_validation_fallback(content)
            
            # Use the real advanced knowledge engine for validation
            validation_result = {
                "entities": [],
                "confidence": 0.0,
                "method": "advanced_validation",
                "enhancement_applied": False
            }
            
            # Basic content validation
            if len(content) < 10:
                validation_result["confidence"] = 0.1
            elif len(content) < 100:
                validation_result["confidence"] = 0.3
            elif len(content) < 1000:
                validation_result["confidence"] = 0.6
            else:
                validation_result["confidence"] = 0.8
            
            # Try to use advanced engine for enhancement if available
            try:
                if hasattr(self.advanced_engine, 'validate_content'):
                    enhanced_result = await self.advanced_engine.validate_content(content, document_type)
                    if enhanced_result:
                        validation_result["confidence"] = max(validation_result["confidence"], 
                                                           enhanced_result.get("confidence", 0.0))
                        validation_result["enhancement_applied"] = True
                        validation_result["enhancement_details"] = enhanced_result
            except Exception as e:
                logger.debug(f"Advanced validation not available: {e}")
                # Continue with basic validation
            
            return validation_result
            
        except Exception as e:
            logger.error(f"Validation failed: {e}")
            return {"entities": [], "confidence": 0.0, "method": "validation_failed"}
    
    def _basic_validation_fallback(self, content: str) -> Dict[str, Any]:
        """Basic validation fallback when advanced engine is not available"""
        confidence = 0.0
        if len(content) < 10:
            confidence = 0.1
        elif len(content) < 100:
            confidence = 0.3
        elif len(content) < 1000:
            confidence = 0.6
        else:
            confidence = 0.8
        
        return {
            "entities": [],
            "confidence": confidence,
            "method": "basic_validation_fallback",
            "enhancement_applied": False
        }
    
    def _merge_processing_results(self, llm_result: Dict[str, Any], 
                                extraction_result: Dict[str, Any],
                                validation_result: Dict[str, Any],
                                content: str, document_type: str, 
                                document_id: str, file_size: int) -> ProcessingResult:
        """Merge results from different processing methods with enhanced AI quality"""
        
        # Combine all entities with quality prioritization
        all_entities = []
        
        # Add LLM entities first (highest quality)
        llm_entities = llm_result.get("entities", [])
        for entity in llm_entities:
            if isinstance(entity, dict):
                entity["source"] = "llm_processing"
                entity["quality_score"] = entity.get("confidence", 0.8)
                all_entities.append(entity)
        
        # Add extraction entities (high quality)
        extraction_entities = extraction_result.get("entities", [])
        for entity in extraction_entities:
            if isinstance(entity, dict):
                entity["source"] = "enhanced_extraction"
                entity["quality_score"] = entity.get("confidence", 0.7)
                all_entities.append(entity)
        
        # Remove duplicates with quality-aware deduplication
        unique_entities = self._deduplicate_entities_quality_aware(all_entities)
        
        # Calculate overall confidence with quality weighting
        llm_confidence = llm_result.get("confidence", 0.0)
        extraction_confidence = extraction_result.get("confidence", 0.0)
        validation_confidence = validation_result.get("confidence", 0.0)
        
        # Weight LLM results higher for better quality
        weighted_confidence = (llm_confidence * 0.5 + 
                             extraction_confidence * 0.3 + 
                             validation_confidence * 0.2)
        
        # Determine processing method and quality level
        methods = [
            llm_result.get("method", "unknown"),
            extraction_result.get("method", "unknown"),
            validation_result.get("method", "unknown")
        ]
        
        # Check if we have high-quality AI processing
        has_llm_processing = llm_result.get("method") != "llm_unavailable" and llm_result.get("method") != "llm_failed"
        has_enhanced_extraction = extraction_result.get("method") == "enhanced_extraction"
        has_advanced_validation = validation_result.get("method") == "advanced_validation"
        
        if has_llm_processing and has_enhanced_extraction:
            processing_method = "ai_enhanced_processing"
            optimization_level = "maximum_quality"
        elif has_llm_processing or has_enhanced_extraction:
            processing_method = "ai_processing"
            optimization_level = "high_quality"
        else:
            processing_method = "basic_processing"
            optimization_level = "standard_quality"
        
        # Generate enhanced content summary
        content_summary = self._generate_enhanced_content_summary(content, unique_entities, document_type)
        
        # Create comprehensive performance metrics
        performance_metrics = {
            "llm_processing_quality": llm_result.get("quality_metrics", {}),
            "extraction_quality": extraction_result.get("extraction_quality", "unknown"),
            "validation_enhancement": validation_result.get("enhancement_applied", False),
            "ai_engine_utilization": {
                "llm_available": has_llm_processing,
                "extraction_available": has_enhanced_extraction,
                "validation_available": has_advanced_validation
            },
            "overall_quality_score": weighted_confidence,
            "processing_efficiency": "high" if weighted_confidence > 0.7 else "medium"
        }
        
        return ProcessingResult(
            document_id=document_id,
            document_type=document_type,
            processing_time=0.0,  # Will be set by caller
            entities_extracted=len(unique_entities),
            confidence_score=weighted_confidence,
            performance_metrics=performance_metrics,
            entities=unique_entities,
            processing_method=processing_method,
            cache_hit=False,  # Will be set by caller
            optimization_level=optimization_level,
            content_summary=content_summary,
            file_size=file_size,
            format_detected=document_type
        )
    
    def _deduplicate_entities(self, entities: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Deduplicate entities (simple approach)"""
        if not entities:
            return []
        
        seen = set()
        unique_entities = []
        
        for entity in entities:
            if isinstance(entity, dict):
                content = entity.get("content", "")
                entity_type = entity.get("entity_type", "")
                key = f"{content[:50]}_{entity_type}"
                
                if key not in seen:
                    unique_entities.append(entity)
                    seen.add(key)
        
        return unique_entities
    
    def _deduplicate_entities_quality_aware(self, entities: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Deduplicate entities with quality-awareness, keeping highest quality versions"""
        if not entities:
            return []
        
        # Sort by quality score descending (highest quality first)
        entities.sort(key=lambda x: x.get('quality_score', 0), reverse=True)
        
        seen_hashes = set()
        unique_entities = []
        
        for entity in entities:
            if isinstance(entity, dict):
                # Create a hash that includes content, type, and category
                content = entity.get('content', '')[:50]
                entity_type = entity.get('entity_type', '')
                category = entity.get('category', '')
                content_hash = f"{content}_{entity_type}_{category}"
                
                if content_hash not in seen_hashes:
                    unique_entities.append(entity)
                    seen_hashes.add(content_hash)
        
        return unique_entities
    
    def _generate_content_summary(self, content: str, entities: List[Dict[str, Any]]) -> str:
        """Generate fast content summary"""
        if not content:
            return "No content available"
        
        # Simple summary based on content length and entities
        word_count = len(content.split())
        entity_count = len(entities)
        
        if word_count < 100:
            summary = f"Short document ({word_count} words)"
        elif word_count < 1000:
            summary = f"Medium document ({word_count} words)"
        else:
            summary = f"Long document ({word_count} words)"
        
        if entity_count > 0:
            summary += f" with {entity_count} extracted entities"
        
        return summary
    
    def _generate_enhanced_content_summary(self, content: str, entities: List[Dict[str, Any]], document_type: str) -> str:
        """Generate enhanced content summary with AI quality metrics."""
        if not content:
            return "No content available"

        word_count = len(content.split())
        entity_count = len(entities)

        summary = f"Document type: {document_type}"
        if word_count < 100:
            summary += f" (Short document, {word_count} words)"
        elif word_count < 1000:
            summary += f" (Medium document, {word_count} words)"
        else:
            summary += f" (Long document, {word_count} words)"

        if entity_count > 0:
            summary += f" with {entity_count} extracted entities"

        # Add AI quality metrics to summary
        if entities:
            avg_quality = sum(e.get('quality_score', 0) for e in entities) / len(entities)
            summary += f" (AI Quality: {avg_quality:.2f})"

        return summary
    
    def _process_video_document(self, file_path: Path, progress_callback: Optional[Callable[[int, int], None]] = None) -> Dict[str, Any]:
        """Compatibility helper for sparse OCR over video frames.
        Samples frames at fixed time intervals and performs fast OCR.
        Returns a dict with 'text' and 'frames_analyzed'.

        progress_callback: optional callable(done_frames, planned_max_frames) to allow UI heartbeat.
        """
        if not self.video_processing_available:
            return {'text': '', 'frames_analyzed': 0}
        
        # Try default backend first
        capture = cv2.VideoCapture(str(file_path))
        if not capture.isOpened():
            # Retry with FFMPEG backend explicitly (helps on some macOS setups)
            try:
                capture = cv2.VideoCapture(str(file_path), cv2.CAP_FFMPEG)
            except Exception:
                capture = capture  # keep as-is
        if not capture.isOpened():
            # As a last resort, try audio-only transcription if available
            if self.audio_processing_available:
                try:
                    import whisper  # type: ignore
                    model = whisper.load_model("base")
                    transcribe = model.transcribe(str(file_path), fp16=False)
                    audio_text = (transcribe.get("text") or "").strip()
                    if audio_text:
                        return {'text': audio_text, 'frames_analyzed': 0, 'source': 'audio_whisper'}
                except Exception:
                    pass
            return {'text': '', 'frames_analyzed': 0, 'warning': 'unable_to_open_video'}
        
        try:
            fps = capture.get(cv2.CAP_PROP_FPS) or 0.0
            total_frames = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
            interval_seconds = getattr(config_manager.processing, 'video_frame_interval_seconds', 5) or 5
            frames_between_samples = max(1, int(round((fps or 1.0) * max(1, interval_seconds))))
            
            # Safety caps for speed: default 50, overridable via env var EXPLAINIUM_MAX_VIDEO_FRAMES
            try:
                env_cap = int(os.getenv("EXPLAINIUM_MAX_VIDEO_FRAMES", "50"))
            except Exception:
                env_cap = 50
            max_frames = min(total_frames, max(5, env_cap))
            frame_index = 0
            frames_analyzed = 0
            collected_lines: List[str] = []
            
            while frame_index < total_frames and frames_analyzed < max_frames:
                capture.set(cv2.CAP_PROP_POS_FRAMES, frame_index)
                success, frame = capture.read()
                if not success or frame is None:
                    frame_index += frames_between_samples
                    continue
                
                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                # Light denoising and thresholding (fast)
                denoised = cv2.fastNlMeansDenoising(gray)
                _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                
                text = self._extract_text_from_frame(thresh) if self.ocr_available else ''
                if text and text.strip():
                    # Keep short, meaningful snippets
                    snippet = " ".join(text.strip().split())
                    if snippet:
                        collected_lines.append(snippet[:200])
                
                frames_analyzed += 1
                if progress_callback:
                    try:
                        progress_callback(frames_analyzed, max_frames)
                    except Exception:
                        pass
                frame_index += frames_between_samples
            
            ocr_text = "\n".join(collected_lines)
            return {'text': ocr_text, 'frames_analyzed': frames_analyzed, 'source': 'frame_ocr'}
        finally:
            capture.release()
    
    def _extract_text_from_frame(self, frame) -> str:
        """Extract text from a video frame using EasyOCR"""
        try:
            import easyocr
            reader = easyocr.Reader(['en'], gpu=False)
            result = reader.readtext(frame)
            
            text_parts = []
            for (bbox, detected_text, confidence) in result:
                if confidence > 0.5:
                    text_parts.append(detected_text)
            
            if text_parts:
                return ' '.join(text_parts)
            return ""
        except Exception as e:
            logger.warning(f"Frame OCR failed: {e}")
            return ""
    
    def _get_file_type(self, extension: str) -> str:
        """Get file type from extension"""
        for file_type, extensions in self.supported_formats.items():
            if extension in extensions:
                return file_type
        return "unknown"
    
    def _generate_cache_key(self, file_path: str) -> str:
        """Generate cache key for file"""
        file_path_obj = Path(file_path)
        file_hash = hashlib.md5(f"{file_path}_{file_path_obj.stat().st_mtime}".encode()).hexdigest()
        return f"{file_path_obj.suffix}_{file_hash}"
    
    def _get_cached_result(self, cache_key: str) -> Optional[ProcessingResult]:
        """Get cached processing result"""
        with self.cache_lock:
            return self.processing_cache.get(cache_key)
    
    def _cache_result(self, cache_key: str, result: ProcessingResult):
        """Cache processing result"""
        with self.cache_lock:
            # Limit cache size for memory efficiency
            if len(self.processing_cache) > 100:
                # Remove oldest entries
                oldest_keys = list(self.processing_cache.keys())[:20]
                for key in oldest_keys:
                    del self.processing_cache[key]
            
            self.processing_cache[cache_key] = result
    
    def _update_performance_stats(self, processing_time: float, format_type: str):
        """Update performance statistics"""
        self.processing_stats["total_documents"] += 1
        self.processing_stats["total_processing_time"] += processing_time
        
        # Calculate running average
        total_docs = self.processing_stats["total_documents"]
        total_time = self.processing_stats["total_processing_time"]
        self.processing_stats["average_processing_time"] = total_time / total_docs
        
        # Track format-specific performance
        if format_type not in self.processing_stats["format_processing_times"]:
            self.processing_stats["format_processing_times"][format_type] = []
        
        self.processing_stats["format_processing_times"][format_type].append(processing_time)
        
        # Keep only last 50 times per format for memory efficiency
        if len(self.processing_stats["format_processing_times"][format_type]) > 50:
            self.processing_stats["format_processing_times"][format_type] = \
                self.processing_stats["format_processing_times"][format_type][-50:]
        
        # Track target performance
        if processing_time <= self.target_processing_time:
            self.processing_stats["performance_target_met"] += 1
        else:
            self.processing_stats["performance_target_missed"] += 1
    
    def _monitor_performance(self, processing_time: float, document_id: str):
        """Monitor and report performance"""
        if processing_time > self.target_processing_time:
            print(f"⚠️ PERFORMANCE WARNING: Document {document_id} took {processing_time:.2f}s "
                  f"(target: {self.target_processing_time}s)")
        elif processing_time > self.performance_warning_threshold:
            print(f"⚠️ PERFORMANCE WARNING: Document {document_id} took {processing_time:.2f}s "
                  f"(approaching target limit)")
        else:
            print(f"✅ Document {document_id} processed in {processing_time:.2f}s "
                  f"(target: {self.target_processing_time}s)")
    
    def get_performance_summary(self) -> Dict[str, Any]:
        """Get comprehensive performance summary"""
        total_docs = self.processing_stats["total_documents"]
        
        if total_docs == 0:
            return {
                "status": "no_documents_processed",
                "performance_optimized": True,
                "target_met": False
            }
        
        cache_hit_rate = (self.processing_stats["cache_hits"] / 
                         (self.processing_stats["cache_hits"] + self.processing_stats["cache_misses"]))
        
        target_success_rate = (self.processing_stats["performance_target_met"] / total_docs)
        
        # Format-specific performance
        format_performance = {}
        for format_type, times in self.processing_stats["format_processing_times"].items():
            if times:
                format_performance[format_type] = {
                    "average_time": sum(times) / len(times),
                    "total_processed": len(times),
                    "target_met_rate": sum(1 for t in times if t <= self.target_processing_time) / len(times)
                }
        
        return {
            "total_documents_processed": total_docs,
            "average_processing_time": self.processing_stats["average_processing_time"],
            "cache_hit_rate": cache_hit_rate,
            "performance_target_success_rate": target_success_rate,
            "performance_target_met_count": self.processing_stats["performance_target_met"],
            "performance_target_missed_count": self.processing_stats["performance_target_missed"],
            "error_count": self.processing_stats["error_count"],
            "format_performance": format_performance,
            "performance_optimized": True,
            "target_met": target_success_rate >= 0.8,  # 80% success rate
            "current_performance_vs_target": {
                "target_time": self.target_processing_time,
                "current_average": self.processing_stats["average_processing_time"],
                "improvement_factor": (10.0 * 60.0) / self.processing_stats["average_processing_time"]  # vs 10 minutes
            }
        }
    
    def optimize_for_m4(self):
        """Apply M4-specific optimizations"""
        print("🔧 Applying M4-specific optimizations...")
        
        # Adjust thread pool for M4 efficiency cores
        self.max_workers = 6  # M4 has 6 efficiency cores
        if self.executor:
            self.executor.shutdown(wait=True)
        self.executor = ThreadPoolExecutor(max_workers=self.max_workers)
        
        # Memory optimization
        self.target_processing_time = 120.0  # 2 minutes
        self.performance_warning_threshold = 90.0  # 1.5 minutes
        
        print(f"✅ M4 optimizations applied: {self.max_workers} workers, "
              f"{self.target_processing_time}s target")
    
    def cleanup(self):
        """Cleanup resources"""
        if self.executor:
            self.executor.shutdown(wait=True)
        
        if self.llm_engine:
            self.llm_engine.cleanup()
        
        if self.extraction_engine:
            self.extraction_engine.cleanup()
        
        # Clear cache
        with self.cache_lock:
            self.processing_cache.clear()
        
        print("🧹 Optimized Document Processor cleaned up")

# Backward compatibility
DocumentProcessor = OptimizedDocumentProcessor